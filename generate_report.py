"""
Coaching Call Discussions — Monthly PowerPoint Report Generator

Produces a professional branded report with:
  - Current month metrics
  - Year-to-date (L12M) totals
  - Month-over-month (MoM) comparisons with delta indicators

Usage:
    python3 generate_report.py                              # HP_SCCAREFIRST, prior month
    python3 generate_report.py ER_SHBP                      # specific customer
    python3 generate_report.py HP_SCCAREFIRST 2026-06-01 2026-06-30
"""
import sys
import os
import logging
from datetime import datetime, date
from dateutil.relativedelta import relativedelta

sys.path.append(os.path.expanduser('~/Documents/dev/automation'))
from db_connect import get_connection
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Config ---
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, 'template.pptx')
DEFAULT_CUSTOMER = 'HP_SCCAREFIRST'

# Brand colors
COLOR_PRIMARY = RGBColor(0x00, 0x4D, 0x3D)    # Dark teal
COLOR_ACCENT = RGBColor(0x00, 0xBF, 0xA5)     # Bright teal
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_LIGHT_GRAY = RGBColor(0xF7, 0xF7, 0xF7)
COLOR_HEADER_BG = RGBColor(0x00, 0x4D, 0x3D)
COLOR_GREEN = RGBColor(0x2E, 0x7D, 0x32)      # Positive delta
COLOR_RED = RGBColor(0xC6, 0x28, 0x28)        # Negative delta
COLOR_MUTED = RGBColor(0x75, 0x75, 0x75)      # Subtle text

FONT_HEADING = 'Roboto Serif 20pt'
FONT_BODY = 'Proxima Nova Rg'

logging.basicConfig(level=logging.INFO, format='%(asctime)s  %(levelname)-8s  %(message)s')
log = logging.getLogger(__name__)

# --- Parse args ---
customer_id = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_CUSTOMER
start_date = sys.argv[2] if len(sys.argv) > 2 else None
end_date = sys.argv[3] if len(sys.argv) > 3 else None

# Default: prior full month
if not start_date:
    today = date.today()
    first_of_this_month = today.replace(day=1)
    last_month_end = first_of_this_month - relativedelta(days=1)
    last_month_start = last_month_end.replace(day=1)
    start_date = last_month_start.strftime('%Y-%m-%d')
    end_date = last_month_end.strftime('%Y-%m-%d')
elif not end_date:
    end_date = date.today().strftime('%Y-%m-%d')

# Compute prior month and trailing 12-month ranges
report_start = datetime.strptime(start_date, '%Y-%m-%d').date()
report_end = datetime.strptime(end_date, '%Y-%m-%d').date()
prior_month_start = (report_start - relativedelta(months=1)).strftime('%Y-%m-%d')
prior_month_end = (report_start - relativedelta(days=1)).strftime('%Y-%m-%d')
t12_start = (report_start - relativedelta(months=11)).strftime('%Y-%m-%d')
t12_end = end_date

# 90-day rolling windows
r90_end = report_end
r90_start = (report_end - relativedelta(days=89)).strftime('%Y-%m-%d')
p90_end = (report_end - relativedelta(days=90)).strftime('%Y-%m-%d')
p90_start = (report_end - relativedelta(days=179)).strftime('%Y-%m-%d')
r90_end_str = r90_end.strftime('%Y-%m-%d')

report_month_label = report_start.strftime('%B %Y')
prior_month_label = (report_start - relativedelta(months=1)).strftime('%B %Y')
t12_label = f"Last 12 Months"


def build_filter(table_alias='', start=None, end=None, date_col='CALL_DATE'):
    """Build WHERE clause fragments for customer and date filtering."""
    s = start or start_date
    e = end or end_date
    parts = []
    prefix = f"{table_alias}." if table_alias else ""
    parts.append(f"UPPER({prefix}CUSTOMERID) = UPPER('{customer_id}')")
    if date_col:
        parts.append(f"{prefix}{date_col} >= '{s}'")
        parts.append(f"{prefix}{date_col} <= '{e}'")
    return ' AND '.join(parts)


# --- Query data ---
log.info(f"Customer: {customer_id}")
log.info(f"Report month: {report_month_label} ({start_date} to {end_date})")
log.info(f"Prior month: {prior_month_label} ({prior_month_start} to {prior_month_end})")
log.info(f"L12M: {t12_start} to {t12_end}")


def query_engagement(conn, start, end):
    return pd.read_sql(f"""
        SELECT
            T.REPORT_TOPIC AS WELLBEING_TOPIC,
            COUNT(DISTINCT T.CURRENTGUID) AS MEMBERS,
            ROUND(COUNT(DISTINCT T.CURRENTGUID) * 100.0
                / NULLIFZERO(TOTALS.TOTAL_MEMBERS), 1) AS PCT_OF_MEMBERS,
            COUNT(DISTINCT CASE WHEN G.GOAL_STATUS = 'Completed' THEN G.MEMBERACTION_ID END) AS COMPLETED_GOALS,
            COUNT(DISTINCT CASE WHEN G.GOAL_STATUS IN ('In Progress','Not Started') THEN G.MEMBERACTION_ID END) AS OPEN_GOALS
        FROM Carefirst_Sandbox.COACHING_CALL_TOPICS T
        CROSS JOIN (
            SELECT COUNT(DISTINCT CURRENTGUID) AS TOTAL_MEMBERS
            FROM Carefirst_Sandbox.COACHING_CALL_TOPICS
            WHERE {build_filter(start=start, end=end)}
        ) TOTALS
        LEFT JOIN Carefirst_Sandbox.COACHING_CALL_GOALS G ON T.CURRENTGUID = G.CURRENTGUID
        WHERE {build_filter('T', start=start, end=end)}
        GROUP BY 1, TOTALS.TOTAL_MEMBERS
        ORDER BY 2 DESC
    """, conn)


def query_goal_dist(conn, start, end):
    return pd.read_sql(f"""
        SELECT
            G.GOAL_STATUS,
            COUNT(*) AS COUNT,
            ROUND(COUNT(*) * 100.0 / NULLIFZERO(SUM(COUNT(*)) OVER()), 1) AS GOAL_PCT
        FROM Carefirst_Sandbox.COACHING_CALL_GOALS G
        JOIN Carefirst_Sandbox.COACHING_CALL_TOPICS T ON G.CURRENTGUID = T.CURRENTGUID
        WHERE G.GOAL_STATUS IN ('Completed','In Progress','Not Started','Withdrawn')
          AND {build_filter('T', start=start, end=end)}
        GROUP BY 1
        ORDER BY CASE G.GOAL_STATUS
            WHEN 'Completed' THEN 1 WHEN 'In Progress' THEN 2
            WHEN 'Not Started' THEN 3 WHEN 'Withdrawn' THEN 4 END
    """, conn)


def query_goal_prog(conn, start, end):
    return pd.read_sql(f"""
        SELECT
            G.GOAL_DOMAIN,
            COUNT(*) AS TOTAL_GOALS,
            SUM(CASE WHEN G.GOAL_STATUS = 'Completed' THEN 1 ELSE 0 END) AS COMPLETED,
            ROUND(SUM(CASE WHEN G.GOAL_STATUS = 'Completed' THEN 1 ELSE 0 END) * 100.0
                / NULLIFZERO(COUNT(*)), 1) AS COMPLETION_RATE
        FROM Carefirst_Sandbox.COACHING_CALL_GOALS G
        JOIN Carefirst_Sandbox.COACHING_CALL_TOPICS T ON G.CURRENTGUID = T.CURRENTGUID
        WHERE {build_filter('T', start=start, end=end)}
        GROUP BY 1
        ORDER BY CASE G.GOAL_DOMAIN
            WHEN 'Gaps in Care' THEN 1 WHEN 'Exercise' THEN 2 WHEN 'Nutrition' THEN 3
            WHEN 'Weight Management' THEN 4 WHEN 'Tobacco Cessation' THEN 5
            WHEN 'Mental/Behavioral Health' THEN 6 WHEN 'Stress Management' THEN 7
            WHEN 'Condition Management' THEN 8 WHEN 'Financial' THEN 9
            WHEN 'Social' THEN 10 WHEN 'Spiritual' THEN 11 ELSE 12 END
    """, conn)


def query_tobacco(conn, start, end):
    return pd.read_sql(f"""
        SELECT 'Tobacco Participants' AS METRIC,
            COUNT(DISTINCT TB.CURRENTGUID)::VARCHAR AS VALUE
        FROM Carefirst_Sandbox.COACHING_CALL_TOBACCO TB
        JOIN Carefirst_Sandbox.COACHING_CALL_TOPICS T ON TB.CURRENTGUID = T.CURRENTGUID
        WHERE {build_filter('T', start=start, end=end)}
        UNION ALL
        SELECT 'Active Tobacco Participants',
            COUNT(DISTINCT G.CURRENTGUID)::VARCHAR
        FROM Carefirst_Sandbox.COACHING_CALL_TOBACCO TB
        JOIN Carefirst_Sandbox.COACHING_CALL_TOPICS T ON TB.CURRENTGUID = T.CURRENTGUID
        JOIN Carefirst_Sandbox.COACHING_CALL_GOALS G
            ON TB.CURRENTGUID = G.CURRENTGUID AND G.GOAL_DOMAIN = 'Tobacco Cessation'
            AND G.GOAL_STATUS IN ('In Progress','Not Started')
        WHERE {build_filter('T', start=start, end=end)}
        UNION ALL
        SELECT 'Goals Completed', COUNT(*)::VARCHAR
        FROM Carefirst_Sandbox.COACHING_CALL_GOALS G
        JOIN Carefirst_Sandbox.COACHING_CALL_TOPICS T ON G.CURRENTGUID = T.CURRENTGUID
        WHERE G.GOAL_DOMAIN = 'Tobacco Cessation' AND G.GOAL_STATUS = 'Completed'
          AND {build_filter('T', start=start, end=end)}
        UNION ALL
        SELECT 'Goals In Progress', COUNT(*)::VARCHAR
        FROM Carefirst_Sandbox.COACHING_CALL_GOALS G
        JOIN Carefirst_Sandbox.COACHING_CALL_TOPICS T ON G.CURRENTGUID = T.CURRENTGUID
        WHERE G.GOAL_DOMAIN = 'Tobacco Cessation' AND G.GOAL_STATUS = 'In Progress'
          AND {build_filter('T', start=start, end=end)}
        UNION ALL
        SELECT 'Completion Rate',
            ROUND(SUM(CASE WHEN G.GOAL_STATUS = 'Completed' THEN 1 ELSE 0 END) * 100.0
                / NULLIFZERO(COUNT(*)), 1)::VARCHAR || '%%'
        FROM Carefirst_Sandbox.COACHING_CALL_GOALS G
        JOIN Carefirst_Sandbox.COACHING_CALL_TOPICS T ON G.CURRENTGUID = T.CURRENTGUID
        WHERE G.GOAL_DOMAIN = 'Tobacco Cessation'
          AND G.GOAL_STATUS IN ('Completed','In Progress','Not Started')
          AND {build_filter('T', start=start, end=end)}
    """, conn)


with get_connection() as conn:
    # Total enrolled members (L12M) — exclude FEP for CareFirst
    cursor = conn.cursor()
    fep_filter = ""  # FEP exclusion handled in monthly_refresh.sql for CareFirst

    cursor.execute(f"""
        SELECT COUNT(DISTINCT CE.CURRENTGUID)
        FROM ENT_WH.COACHING_ENROLLMENT_MODEL CE
        WHERE UPPER(CE.LEVEL_NAME) = 'ENROLLED' AND UPPER(CE.OFFERING_STATUS) = 'ENROLLED'
          AND CE.BIEFFECTIVEDATE <= '{t12_end}'
          AND CE.BIENDDATE >= '{t12_start}'
          AND UPPER(CE.CUSTOMERID) = UPPER('{customer_id}')
          {fep_filter}
    """)
    total_enrolled_t12 = cursor.fetchone()[0]

    cursor.execute(f"""
        SELECT COUNT(DISTINCT CE.CURRENTGUID)
        FROM ENT_WH.COACHING_ENROLLMENT_MODEL CE
        WHERE UPPER(CE.LEVEL_NAME) = 'ENROLLED' AND UPPER(CE.OFFERING_STATUS) = 'ENROLLED'
          AND CE.BIEFFECTIVEDATE <= '{end_date}'
          AND CE.BIENDDATE >= '{start_date}'
          AND UPPER(CE.CUSTOMERID) = UPPER('{customer_id}')
          {fep_filter}
    """)
    total_enrolled_month = cursor.fetchone()[0]

    cursor.execute(f"""
        SELECT COUNT(DISTINCT CE.CURRENTGUID)
        FROM ENT_WH.COACHING_ENROLLMENT_MODEL CE
        WHERE UPPER(CE.LEVEL_NAME) = 'ENROLLED' AND UPPER(CE.OFFERING_STATUS) = 'ENROLLED'
          AND CE.BIEFFECTIVEDATE <= '{prior_month_end}'
          AND CE.BIENDDATE >= '{prior_month_start}'
          AND UPPER(CE.CUSTOMERID) = UPPER('{customer_id}')
          {fep_filter}
    """)
    total_enrolled_prior = cursor.fetchone()[0]

    # 90-day enrollment counts
    cursor.execute(f"""
        SELECT COUNT(DISTINCT CE.CURRENTGUID)
        FROM ENT_WH.COACHING_ENROLLMENT_MODEL CE
        WHERE UPPER(CE.LEVEL_NAME) = 'ENROLLED' AND UPPER(CE.OFFERING_STATUS) = 'ENROLLED'
          AND CE.BIEFFECTIVEDATE <= '{r90_end_str}'
          AND CE.BIENDDATE >= '{r90_start}'
          AND UPPER(CE.CUSTOMERID) = UPPER('{customer_id}')
          {fep_filter}
    """)
    enrolled_r90 = cursor.fetchone()[0]

    cursor.execute(f"""
        SELECT COUNT(DISTINCT CE.CURRENTGUID)
        FROM ENT_WH.COACHING_ENROLLMENT_MODEL CE
        WHERE UPPER(CE.LEVEL_NAME) = 'ENROLLED' AND UPPER(CE.OFFERING_STATUS) = 'ENROLLED'
          AND CE.BIEFFECTIVEDATE <= '{p90_end}'
          AND CE.BIENDDATE >= '{p90_start}'
          AND UPPER(CE.CUSTOMERID) = UPPER('{customer_id}')
          {fep_filter}
    """)
    enrolled_p90 = cursor.fetchone()[0]

    enrolled_r90_avg = round(enrolled_r90 / 3)
    enrolled_p90_avg = round(enrolled_p90 / 3)

    # Current month
    df_engagement = query_engagement(conn, start_date, end_date)
    df_goal_dist = query_goal_dist(conn, start_date, end_date)
    df_goal_prog = query_goal_prog(conn, start_date, end_date)
    df_tobacco = query_tobacco(conn, start_date, end_date)

    # Prior month
    df_engagement_prior = query_engagement(conn, prior_month_start, prior_month_end)
    df_goal_dist_prior = query_goal_dist(conn, prior_month_start, prior_month_end)
    df_tobacco_prior = query_tobacco(conn, prior_month_start, prior_month_end)

    # L12M
    df_engagement_t12 = query_engagement(conn, t12_start, t12_end)
    df_goal_dist_t12 = query_goal_dist(conn, t12_start, t12_end)
    df_goal_prog_t12 = query_goal_prog(conn, t12_start, t12_end)
    df_tobacco_t12 = query_tobacco(conn, t12_start, t12_end)

    # 90-day rolling windows
    df_engagement_r90 = query_engagement(conn, r90_start, r90_end_str)
    df_engagement_p90 = query_engagement(conn, p90_start, p90_end)
    df_goal_dist_r90 = query_goal_dist(conn, r90_start, r90_end_str)
    df_goal_dist_p90 = query_goal_dist(conn, p90_start, p90_end)
    df_goal_prog_r90 = query_goal_prog(conn, r90_start, r90_end_str)
    df_goal_prog_p90 = query_goal_prog(conn, p90_start, p90_end)

    # Classic summary tables (L12M window for the overview slide)
    df_topics_classic = pd.read_sql(f"""
        SELECT REPORT_TOPIC AS CALL_TOPIC,
            COUNT(DISTINCT CURRENTGUID) AS MEMBERS,
            ROUND(COUNT(DISTINCT CURRENTGUID)*100.0/NULLIFZERO(TOTALS.TOTAL_MEMBERS),2) AS PCT_MEMBERS,
            COUNT(*) AS CALLS,
            ROUND(COUNT(*)*100.0/NULLIFZERO(TOTALS.TOTAL_CALLS),2) AS PCT_OF_CALLS
        FROM Carefirst_Sandbox.COACHING_CALL_TOPICS T
        CROSS JOIN (
            SELECT COUNT(DISTINCT CURRENTGUID) AS TOTAL_MEMBERS, COUNT(*) AS TOTAL_CALLS
            FROM Carefirst_Sandbox.COACHING_CALL_TOPICS WHERE {build_filter(start=t12_start, end=t12_end)}
        ) TOTALS
        WHERE {build_filter('T', start=t12_start, end=t12_end)}
        GROUP BY 1, TOTALS.TOTAL_MEMBERS, TOTALS.TOTAL_CALLS
        ORDER BY 4 DESC
    """, conn)

    df_tobacco_classic = pd.read_sql(f"""
        SELECT 'Yes' AS TOBACCO_DISCUSSED, COUNT(DISTINCT T.CURRENTGUID) AS MEMBERS,
            ROUND(COUNT(DISTINCT T.CURRENTGUID)*100.0/NULLIFZERO(TOTALS.TOTAL_MEMBERS),2) AS PCT_MEMBERS,
            COUNT(*) AS CALLS,
            ROUND(COUNT(*)*100.0/NULLIFZERO(TOTALS.TOTAL_CALLS),2) AS PCT_OF_CALLS
        FROM Carefirst_Sandbox.COACHING_CALL_TOPICS T
        JOIN Carefirst_Sandbox.COACHING_CALL_TOBACCO TB ON T.CURRENTGUID = TB.CURRENTGUID
        CROSS JOIN (SELECT COUNT(DISTINCT CURRENTGUID) AS TOTAL_MEMBERS, COUNT(*) AS TOTAL_CALLS
            FROM Carefirst_Sandbox.COACHING_CALL_TOPICS WHERE {build_filter(start=t12_start, end=t12_end)}) TOTALS
        WHERE {build_filter('T', start=t12_start, end=t12_end)}
        GROUP BY 1, TOTALS.TOTAL_MEMBERS, TOTALS.TOTAL_CALLS
        UNION ALL
        SELECT 'No', COUNT(DISTINCT T.CURRENTGUID),
            ROUND(COUNT(DISTINCT T.CURRENTGUID)*100.0/NULLIFZERO(TOTALS.TOTAL_MEMBERS),2),
            COUNT(*), ROUND(COUNT(*)*100.0/NULLIFZERO(TOTALS.TOTAL_CALLS),2)
        FROM Carefirst_Sandbox.COACHING_CALL_TOPICS T
        CROSS JOIN (SELECT COUNT(DISTINCT CURRENTGUID) AS TOTAL_MEMBERS, COUNT(*) AS TOTAL_CALLS
            FROM Carefirst_Sandbox.COACHING_CALL_TOPICS WHERE {build_filter(start=t12_start, end=t12_end)}) TOTALS
        WHERE {build_filter('T', start=t12_start, end=t12_end)}
          AND T.CURRENTGUID NOT IN (SELECT CURRENTGUID FROM Carefirst_Sandbox.COACHING_CALL_TOBACCO)
        GROUP BY 1, TOTALS.TOTAL_MEMBERS, TOTALS.TOTAL_CALLS
    """, conn)

    df_goal_domain_status = pd.read_sql(f"""
        SELECT G.GOAL_DOMAIN,
            SUM(CASE WHEN G.GOAL_STATUS='Completed' THEN 1 ELSE 0 END) AS COMPLETED,
            ROUND(SUM(CASE WHEN G.GOAL_STATUS='Completed' THEN 1 ELSE 0 END)*100.0/NULLIFZERO(COUNT(*)),1) AS COMP_PCT,
            SUM(CASE WHEN G.GOAL_STATUS='In Progress' THEN 1 ELSE 0 END) AS IN_PROGRESS,
            ROUND(SUM(CASE WHEN G.GOAL_STATUS='In Progress' THEN 1 ELSE 0 END)*100.0/NULLIFZERO(COUNT(*)),1) AS IP_PCT,
            SUM(CASE WHEN G.GOAL_STATUS='Withdrawn' THEN 1 ELSE 0 END) AS WITHDRAWN,
            ROUND(SUM(CASE WHEN G.GOAL_STATUS='Withdrawn' THEN 1 ELSE 0 END)*100.0/NULLIFZERO(COUNT(*)),1) AS WD_PCT
        FROM Carefirst_Sandbox.COACHING_CALL_GOALS G
        JOIN Carefirst_Sandbox.COACHING_CALL_TOPICS T ON G.CURRENTGUID = T.CURRENTGUID
        WHERE {build_filter('T', start=t12_start, end=t12_end)}
        GROUP BY 1
        ORDER BY CASE G.GOAL_DOMAIN WHEN 'Gaps in Care' THEN 1 WHEN 'Exercise' THEN 2 WHEN 'Nutrition' THEN 3
            WHEN 'Weight Management' THEN 4 WHEN 'Tobacco Cessation' THEN 5 WHEN 'Mental/Behavioral Health' THEN 6
            WHEN 'Stress Management' THEN 7 WHEN 'Condition Management' THEN 8 ELSE 9 END
    """, conn)

log.info(f"Data loaded: {len(df_engagement)} topics, {len(df_goal_prog)} domains")


# --- PowerPoint helpers ---
def set_cell(cell, text, font_size=9, bold=False, color=COLOR_DARK, align=PP_ALIGN.LEFT, fill=None):
    """Format a single table cell."""
    cell.text = str(text) if text is not None else ''
    for para in cell.text_frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def add_table(slide, df, left, top, width, height, pct_cols=None, first_col_width=None):
    """Add a branded table with header styling and alternating rows."""
    rows = len(df) + 1
    cols = len(df.columns)
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table

    # Column widths
    if first_col_width and cols > 1:
        tbl.columns[0].width = first_col_width
        remaining = (width - first_col_width) // (cols - 1)
        for i in range(1, cols):
            tbl.columns[i].width = remaining
    else:
        cw = width // cols
        for i in range(cols):
            tbl.columns[i].width = cw

    # Header
    for c, col_name in enumerate(df.columns):
        display = col_name.replace('_', ' ')
        set_cell(tbl.cell(0, c), display, font_size=8, bold=True,
                 color=COLOR_WHITE, align=PP_ALIGN.CENTER, fill=COLOR_HEADER_BG)

    # Data
    for r_idx, (_, row) in enumerate(df.iterrows()):
        for c, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c)
            if pct_cols and df.columns[c] in pct_cols and val is not None:
                display_val = f"{val}%"
            elif isinstance(val, (int,)):
                display_val = f"{val:,}"
            elif isinstance(val, float) and val == int(val):
                display_val = f"{int(val):,}"
            else:
                display_val = val
            # First column left-aligned, all others centered
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            fill = COLOR_LIGHT_GRAY if r_idx % 2 == 0 else None
            set_cell(cell, display_val, font_size=9, color=COLOR_DARK, align=align, fill=fill)

    return tbl_shape


def add_title(slide, text, left, top, size=14):
    """Section title."""
    tb = slide.shapes.add_textbox(left, top, Inches(9), Pt(size * 2))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.name = FONT_BODY
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    return tb


def add_kpi_box(slide, label, value, curr_value=None, prior_value=None, delta=None, delta_label="vs prior month", r90_avg=None, p90_avg=None, left=Inches(0), top=Inches(0), width=Inches(2.2), height=Inches(1.1)):
    """Add a KPI box with monthly comparison and 90-day rolling averages."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True

    # Value (large)
    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.name = FONT_BODY
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Label
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.name = FONT_BODY
    p2.font.size = Pt(9)
    p2.font.color.rgb = COLOR_MUTED
    p2.alignment = PP_ALIGN.CENTER

    # MoM detail line (e.g., "Jun: 3,903 | May: 3,805")
    if prior_value is not None:
        p_mom = tf.add_paragraph()
        short_curr = report_month_label.split()[0][:3]
        short_prior = prior_month_label.split()[0][:3]
        curr_display = curr_value if curr_value is not None else value
        p_mom.text = f"{short_curr}: {curr_display}  |  {short_prior}: {prior_value}"
        p_mom.font.name = FONT_BODY
        p_mom.font.size = Pt(7)
        p_mom.font.color.rgb = COLOR_MUTED
        p_mom.alignment = PP_ALIGN.CENTER

    # Delta indicator
    if delta is not None and delta != 0:
        p3 = tf.add_paragraph()
        arrow = "▲" if delta > 0 else "▼"
        p3.text = f"{arrow} {abs(delta):,} {delta_label}"
        p3.font.name = FONT_BODY
        p3.font.size = Pt(8)
        p3.font.color.rgb = COLOR_GREEN if delta > 0 else COLOR_RED
        p3.alignment = PP_ALIGN.CENTER

    # 90-day rolling average line
    if r90_avg is not None and p90_avg is not None:
        p4 = tf.add_paragraph()
        p4.text = f"{r90_avg}  |  {p90_avg}"
        p4.font.name = FONT_BODY
        p4.font.size = Pt(7)
        p4.font.color.rgb = COLOR_MUTED
        p4.alignment = PP_ALIGN.CENTER

    return tb


def safe_int(val):
    """Convert a value to int, handling strings and None."""
    if val is None:
        return 0
    try:
        return int(float(str(val).replace('%', '').replace(',', '')))
    except (ValueError, TypeError):
        return 0


# --- Compute KPIs ---
current_members = int(df_engagement['MEMBERS'].sum()) if len(df_engagement) > 0 else 0
prior_members = int(df_engagement_prior['MEMBERS'].sum()) if len(df_engagement_prior) > 0 else 0
t12_members = int(df_engagement_t12['MEMBERS'].sum()) if len(df_engagement_t12) > 0 else 0

current_completed = int(df_goal_dist[df_goal_dist['GOAL_STATUS'] == 'Completed']['COUNT'].sum()) if len(df_goal_dist) > 0 else 0
prior_completed = int(df_goal_dist_prior[df_goal_dist_prior['GOAL_STATUS'] == 'Completed']['COUNT'].sum()) if len(df_goal_dist_prior) > 0 else 0

current_calls = len(df_engagement) if len(df_engagement) > 0 else 0  # topics count

# Tobacco KPIs
def clean_tob_value(val):
    """Clean tobacco values — round percentages to 1 decimal."""
    if val is None:
        return '0'
    val = str(val).strip()
    if '%' in val:
        try:
            num = float(val.replace('%', '').replace('%%', ''))
            return f"{num:.1f}%"
        except ValueError:
            return val
    return val

tob_current = {k: clean_tob_value(v) for k, v in (df_tobacco.set_index('METRIC')['VALUE'].to_dict().items())} if len(df_tobacco) > 0 else {}
tob_prior = {k: clean_tob_value(v) for k, v in (df_tobacco_prior.set_index('METRIC')['VALUE'].to_dict().items())} if len(df_tobacco_prior) > 0 else {}
tob_t12 = {k: clean_tob_value(v) for k, v in (df_tobacco_t12.set_index('METRIC')['VALUE'].to_dict().items())} if len(df_tobacco_t12) > 0 else {}


# --- Build presentation ---
prs = Presentation(TEMPLATE_PATH)

# Remove template slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[0]

BLANK = prs.slide_layouts[7]

# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(BLANK)
tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(8.5), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Coaching Call Discussions"
p.font.name = FONT_HEADING
p.font.size = Pt(32)
p.font.color.rgb = COLOR_PRIMARY
p2 = tf.add_paragraph()
p2.text = "Monthly Report"
p2.font.name = FONT_BODY
p2.font.size = Pt(18)
p2.font.color.rgb = COLOR_ACCENT
p2.space_before = Pt(4)

tb2 = slide.shapes.add_textbox(Inches(0.75), Inches(1.9), Inches(8), Inches(0.6))
tf2 = tb2.text_frame
p3 = tf2.paragraphs[0]
p3.text = f"{customer_id}  |  {report_month_label}"
p3.font.name = FONT_BODY
p3.font.size = Pt(13)
p3.font.color.rgb = COLOR_DARK
p4 = tf2.add_paragraph()
p4.text = f"Trailing 12 Months: {t12_start} to {t12_end}  |  Monthly Detail: {report_month_label}"
p4.font.name = FONT_BODY
p4.font.size = Pt(9)
p4.font.color.rgb = COLOR_MUTED

# ===== SLIDE 2: Executive Summary KPIs =====
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Executive Summary", Inches(0.4), Inches(0.2), size=16)

# Subtitle
tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.55), Inches(8), Inches(0.3))
p = tb.text_frame.paragraphs[0]
p.text = f"Last 12 Months ({t12_start[:4]})  |  Monthly Detail: {report_month_label} vs {prior_month_label}"
p.font.name = FONT_BODY
p.font.size = Pt(10)
p.font.color.rgb = COLOR_MUTED

# KPI row 1 — 4 L12M KPIs (with 90d avg subscript)
x_start = Inches(0.2)
y_pos = Inches(0.9)
kpi_width = Inches(2.3)
spacing = Inches(2.45)

# 90-day rolling averages (per month)
r90_members = int(df_engagement_r90['MEMBERS'].sum()) if len(df_engagement_r90) > 0 else 0
p90_members = int(df_engagement_p90['MEMBERS'].sum()) if len(df_engagement_p90) > 0 else 0
r90_completed = int(df_goal_dist_r90[df_goal_dist_r90['GOAL_STATUS'] == 'Completed']['COUNT'].sum()) if len(df_goal_dist_r90) > 0 else 0
p90_completed = int(df_goal_dist_p90[df_goal_dist_p90['GOAL_STATUS'] == 'Completed']['COUNT'].sum()) if len(df_goal_dist_p90) > 0 else 0
r90_avg_members = round(r90_members / 3)
p90_avg_members = round(p90_members / 3)
r90_avg_completed = round(r90_completed / 3)
p90_avg_completed = round(p90_completed / 3)
tob_current_count = tob_current.get('Tobacco Participants', '0')
tob_t12_count = tob_t12.get('Tobacco Participants', '0')
tob_r90 = safe_int(tob_current_count)
tob_p90 = safe_int(tob_prior.get('Tobacco Participants', 0))

add_kpi_box(slide, "L12M Total Enrolled", f"{total_enrolled_t12:,}",
            r90_avg=f"90d Avg: {enrolled_r90_avg:,}", p90_avg=f"Prior 90d: {enrolled_p90_avg:,}",
            delta=enrolled_r90_avg - enrolled_p90_avg,
            delta_label="vs prior 90d",
            left=x_start, top=y_pos, width=kpi_width, height=Inches(1.2))

add_kpi_box(slide, "L12M Members Coached", f"{t12_members:,}",
            r90_avg=f"90d Avg: {r90_avg_members:,}", p90_avg=f"Prior 90d: {p90_avg_members:,}",
            delta=r90_avg_members - p90_avg_members,
            delta_label="vs prior 90d",
            left=x_start + spacing, top=y_pos, width=kpi_width, height=Inches(1.2))

t12_completed = int(df_goal_dist_t12[df_goal_dist_t12['GOAL_STATUS'] == 'Completed']['COUNT'].sum()) if len(df_goal_dist_t12) > 0 else 0
add_kpi_box(slide, "L12M Goals Completed", f"{t12_completed:,}",
            r90_avg=f"90d Avg: {r90_avg_completed:,}", p90_avg=f"Prior 90d: {p90_avg_completed:,}",
            delta=r90_avg_completed - p90_avg_completed,
            delta_label="vs prior 90d",
            left=x_start + spacing * 2, top=y_pos, width=kpi_width, height=Inches(1.2))

add_kpi_box(slide, "L12M Tobacco Participants", tob_t12_count,
            r90_avg=f"90d Avg: {tob_r90:,}", p90_avg=f"Prior 90d: {tob_p90:,}",
            delta=tob_r90 - tob_p90,
            delta_label="vs prior 90d",
            left=x_start + spacing * 3, top=y_pos, width=kpi_width, height=Inches(1.2))

# KPI row 2 — 4 Monthly KPIs (same structure: value, label, prior month line, delta)
y_pos2 = Inches(2.3)

add_kpi_box(slide, f"{report_month_label} Enrolled", f"{total_enrolled_month:,}",
            r90_avg=f"{total_enrolled_month:,}", p90_avg=f"{prior_month_label}: {total_enrolled_prior:,}",
            delta=total_enrolled_month - total_enrolled_prior,
            delta_label="vs prior month",
            left=x_start, top=y_pos2, width=kpi_width, height=Inches(1.2))

add_kpi_box(slide, f"{report_month_label} Members", f"{current_members:,}",
            r90_avg=f"{current_members:,}", p90_avg=f"{prior_month_label}: {prior_members:,}",
            delta=current_members - prior_members,
            delta_label="vs prior month",
            left=x_start + spacing, top=y_pos2, width=kpi_width, height=Inches(1.2))

add_kpi_box(slide, f"{report_month_label} Goals Completed", f"{current_completed:,}",
            r90_avg=f"{current_completed:,}", p90_avg=f"{prior_month_label}: {prior_completed:,}",
            delta=current_completed - prior_completed,
            delta_label="vs prior month",
            left=x_start + spacing * 2, top=y_pos2, width=kpi_width, height=Inches(1.2))

add_kpi_box(slide, f"{report_month_label} Tobacco", tob_current_count,
            r90_avg=f"{tob_current_count}", p90_avg=f"{prior_month_label}: {tob_prior.get('Tobacco Participants', '0')}",
            delta=safe_int(tob_current_count) - safe_int(tob_prior.get('Tobacco Participants', 0)),
            delta_label="vs prior month",
            left=x_start + spacing * 3, top=y_pos2, width=kpi_width, height=Inches(1.2))

# Goal Status Distribution
add_title(slide, "Goal Status Distribution", Inches(0.4), Inches(3.7))

# Build combined DataFrame — L12M first
statuses = ['Completed', 'In Progress', 'Not Started', 'Withdrawn']
combined_rows = []
for status in statuses:
    curr_row = df_goal_dist[df_goal_dist['GOAL_STATUS'] == status]
    prior_row = df_goal_dist_prior[df_goal_dist_prior['GOAL_STATUS'] == status]
    ytd_row = df_goal_dist_t12[df_goal_dist_t12['GOAL_STATUS'] == status]
    curr_count = int(curr_row['COUNT'].iloc[0]) if len(curr_row) > 0 else 0
    curr_pct = float(curr_row['GOAL_PCT'].iloc[0]) if len(curr_row) > 0 else 0.0
    prior_count = int(prior_row['COUNT'].iloc[0]) if len(prior_row) > 0 else 0
    ytd_count = int(ytd_row['COUNT'].iloc[0]) if len(ytd_row) > 0 else 0
    ytd_pct = float(ytd_row['GOAL_PCT'].iloc[0]) if len(ytd_row) > 0 else 0.0
    mom_change = curr_count - prior_count
    combined_rows.append({
        'Goal Status': status,
        'L12M': f"{ytd_count:,}",
        'L12M %': f"{ytd_pct:.1f}%",
        f'{report_month_label}': f"{curr_count:,}",
        f'{prior_month_label}': f"{prior_count:,}",
        'MoM Change': f"+{mom_change:,}" if mom_change > 0 else f"{mom_change:,}",
    })

df_goal_dist_combined = pd.DataFrame(combined_rows)
add_table(slide, df_goal_dist_combined,
          left=Inches(0.4), top=Inches(4.1),
          width=Inches(9.2), height=Inches(2.0),
          first_col_width=Inches(1.6))

# ===== SLIDE 3: Coaching Engagement by Wellbeing Topic =====
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Coaching Engagement by Wellbeing Topic", Inches(0.4), Inches(0.2))

# Build combined engagement table with MoM + L12M + 90-day
engagement_combined_rows = []
all_topics = list(df_engagement['WELLBEING_TOPIC'].unique())
for topic in all_topics:
    curr = df_engagement[df_engagement['WELLBEING_TOPIC'] == topic]
    prior = df_engagement_prior[df_engagement_prior['WELLBEING_TOPIC'] == topic] if len(df_engagement_prior) > 0 else pd.DataFrame()
    ytd = df_engagement_t12[df_engagement_t12['WELLBEING_TOPIC'] == topic] if len(df_engagement_t12) > 0 else pd.DataFrame()
    r90 = df_engagement_r90[df_engagement_r90['WELLBEING_TOPIC'] == topic] if len(df_engagement_r90) > 0 else pd.DataFrame()
    p90 = df_engagement_p90[df_engagement_p90['WELLBEING_TOPIC'] == topic] if len(df_engagement_p90) > 0 else pd.DataFrame()

    curr_members = int(curr['MEMBERS'].iloc[0]) if len(curr) > 0 else 0
    prior_members_val = int(prior['MEMBERS'].iloc[0]) if len(prior) > 0 else 0
    t12_members_val = int(ytd['MEMBERS'].iloc[0]) if len(ytd) > 0 else 0
    r90_mem = int(r90['MEMBERS'].iloc[0]) if len(r90) > 0 else 0
    p90_mem = int(p90['MEMBERS'].iloc[0]) if len(p90) > 0 else 0
    mom_delta = curr_members - prior_members_val
    r90_delta = r90_mem - p90_mem
    curr_pct = float(curr['PCT_OF_MEMBERS'].iloc[0]) if len(curr) > 0 else 0.0

    engagement_combined_rows.append({
        'Wellbeing Topic': topic,
        'L12M': f"{t12_members_val:,}",
        f'{report_month_label}': f"{curr_members:,}",
        f'{prior_month_label}': f"{prior_members_val:,}",
        'MoM \u0394': f"+{mom_delta:,}" if mom_delta > 0 else f"{mom_delta:,}",
        '90d Avg': f"{round(r90_mem/3):,}",
        'Prior 90d Avg': f"{round(p90_mem/3):,}",
        '90d \u0394': f"+{round(r90_delta/3):,}" if r90_delta > 0 else f"{round(r90_delta/3):,}",
    })

df_engagement_combined = pd.DataFrame(engagement_combined_rows)
add_table(slide, df_engagement_combined,
          left=Inches(0.2), top=Inches(0.6),
          width=Inches(9.6), height=Inches(5.5),
          first_col_width=Inches(1.8))


# ===== SLIDE 4: Goal Progression by Domain =====
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Goal Progression by Domain", Inches(0.4), Inches(0.2))

# Build combined goal progression with MoM + L12M + 90d
prog_combined_rows = []
all_domains = list(df_goal_prog['GOAL_DOMAIN'].unique()) if len(df_goal_prog) > 0 else []
for domain in all_domains:
    curr = df_goal_prog[df_goal_prog['GOAL_DOMAIN'] == domain]
    ytd = df_goal_prog_t12[df_goal_prog_t12['GOAL_DOMAIN'] == domain] if len(df_goal_prog_t12) > 0 else pd.DataFrame()
    r90 = df_goal_prog_r90[df_goal_prog_r90['GOAL_DOMAIN'] == domain] if len(df_goal_prog_r90) > 0 else pd.DataFrame()
    p90 = df_goal_prog_p90[df_goal_prog_p90['GOAL_DOMAIN'] == domain] if len(df_goal_prog_p90) > 0 else pd.DataFrame()

    curr_total = int(curr['TOTAL_GOALS'].iloc[0]) if len(curr) > 0 else 0
    curr_completed = int(curr['COMPLETED'].iloc[0]) if len(curr) > 0 else 0
    curr_rate = float(curr['COMPLETION_RATE'].iloc[0]) if len(curr) > 0 else 0.0
    ytd_total = int(ytd['TOTAL_GOALS'].iloc[0]) if len(ytd) > 0 else 0
    t12_completed = int(ytd['COMPLETED'].iloc[0]) if len(ytd) > 0 else 0
    ytd_rate = float(ytd['COMPLETION_RATE'].iloc[0]) if len(ytd) > 0 else 0.0
    r90_total = int(r90['TOTAL_GOALS'].iloc[0]) if len(r90) > 0 else 0
    p90_total = int(p90['TOTAL_GOALS'].iloc[0]) if len(p90) > 0 else 0
    r90_avg = round(r90_total / 3)
    p90_avg = round(p90_total / 3)
    r90_d = r90_avg - p90_avg

    prog_combined_rows.append({
        'Goal Domain': domain,
        'L12M Goals': f"{ytd_total:,}",
        'L12M Completed': f"{t12_completed:,}",
        'L12M Rate': f"{ytd_rate:.1f}%",
        '90D Avg': f"{r90_avg:,}",
        'Prior 90D Avg': f"{p90_avg:,}",
        '90D \u0394': f"+{r90_d:,}" if r90_d > 0 else f"{r90_d:,}",
    })

df_prog_combined = pd.DataFrame(prog_combined_rows)
add_table(slide, df_prog_combined,
          left=Inches(0.2), top=Inches(0.6),
          width=Inches(9.6), height=Inches(5.5),
          first_col_width=Inches(2.2))


# ===== SLIDE 5: Tobacco Coaching Focus =====
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Tobacco Coaching Focus", Inches(0.4), Inches(0.2))

# Build combined tobacco table with month names
tob_metrics = ['Tobacco Participants\u00b9', 'Active Tobacco Participants\u00b2', 'Goals Completed', 'Goals In Progress', 'Completion Rate']
tobacco_combined_rows = []
for metric in tob_metrics:
    curr_val = tob_current.get(metric.rstrip('\u00b9\u00b2'), '0')
    prior_val = tob_prior.get(metric.rstrip('\u00b9\u00b2'), '0')
    ytd_val = tob_t12.get(metric.rstrip('\u00b9\u00b2'), '0')

    # Compute MoM change and 90d (skip for Completion Rate)
    if metric != 'Completion Rate':
        c = safe_int(curr_val)
        p = safe_int(prior_val)
        mom = c - p
        mom_str = f"+{mom}" if mom > 0 else str(mom)
        r90_str = str(c)
        p90_str = str(p)
        d90 = c - p
        d90_str = f"+{d90}" if d90 > 0 else str(d90)
    else:
        mom_str = '—'
        r90_str = '—'
        p90_str = '—'
        d90_str = '—'

    tobacco_combined_rows.append({
        'Metric': metric,
        'L12M': ytd_val,
        f'{report_month_label}': curr_val,
        f'{prior_month_label}': prior_val,
        'MoM Change': mom_str,
        '90D Avg': r90_str,
        'Prior 90D Avg': p90_str,
        '90D Δ': d90_str,
    })

df_tobacco_combined = pd.DataFrame(tobacco_combined_rows)
add_table(slide, df_tobacco_combined,
          left=Inches(0.3), top=Inches(0.6),
          width=Inches(9.4), height=Inches(2.5),
          first_col_width=Inches(2.8))

# Definitions and gap analysis
tob_participants = safe_int(tob_current.get('Tobacco Participants', 0))
tob_active = safe_int(tob_current.get('Active Tobacco Participants', 0))
tob_gap = tob_participants - tob_active
gap_pct = round(tob_gap * 100.0 / max(tob_participants, 1), 1)

tb = slide.shapes.add_textbox(Inches(0.3), Inches(3.3), Inches(9.4), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "Definitions"
p.font.name = FONT_BODY
p.font.size = Pt(10)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

p2 = tf.add_paragraph()
p2.text = "\u00b9 Tobacco Participants: Members who discussed tobacco during a coaching call (topic form selection or Tobacco call type)"
p2.font.name = FONT_BODY
p2.font.size = Pt(9)
p2.font.color.rgb = COLOR_DARK
p2.space_before = Pt(4)

p3 = tf.add_paragraph()
p3.text = "\u00b2 Active Tobacco Participants: Subset with a formal Tobacco Cessation goal in the system (status: In Progress or Not Started)"
p3.font.name = FONT_BODY
p3.font.size = Pt(9)
p3.font.color.rgb = COLOR_DARK
p3.space_before = Pt(2)

p4 = tf.add_paragraph()
p4.text = ""
p4.space_before = Pt(10)

p5 = tf.add_paragraph()
p5.text = "Goal Coverage Gap"
p5.font.name = FONT_BODY
p5.font.size = Pt(10)
p5.font.bold = True
p5.font.color.rgb = COLOR_PRIMARY
p5.space_before = Pt(6)

p6 = tf.add_paragraph()
p6.text = f"{tob_gap} of {tob_participants} tobacco participants ({gap_pct}%) discussed tobacco in coaching but do not have a formal Tobacco Cessation goal set. These members may have goals in other domains (Exercise, Condition Management, etc.) but no tobacco-specific goal is being tracked."
p6.font.name = FONT_BODY
p6.font.size = Pt(9)
p6.font.color.rgb = COLOR_DARK
p6.space_before = Pt(4)


# ===== SLIDE 6: Data Dictionary (Page 1) =====
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Data Dictionary", Inches(0.4), Inches(0.15), size=16)

line = slide.shapes.add_connector(1, Inches(0.4), Inches(0.5), Inches(9.6), Inches(0.5))
line.line.color.rgb = COLOR_ACCENT
line.line.width = Pt(3)

dd1 = slide.shapes.add_textbox(Inches(0.4), Inches(0.6), Inches(9.2), Inches(7.0))
dd1_tf = dd1.text_frame
dd1_tf.word_wrap = True

def dd_section(tf, title):
    p = tf.add_paragraph()
    p.text = title
    p.font.name = FONT_BODY
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_before = Pt(12)

def dd_entry(tf, term, definition):
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    run1 = p.add_run()
    run1.text = term
    run1.font.name = FONT_BODY
    run1.font.size = Pt(8)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_DARK
    run2 = p.add_run()
    run2.text = f" – {definition}"
    run2.font.name = FONT_BODY
    run2.font.size = Pt(8)
    run2.font.color.rgb = COLOR_MUTED

# Page 1: Report Overview + Coaching Engagement + Goals
p0 = dd1_tf.paragraphs[0]
p0.text = "About This Report"
p0.font.name = FONT_BODY
p0.font.size = Pt(11)
p0.font.bold = True
p0.font.color.rgb = COLOR_PRIMARY

dd_entry(dd1_tf, "Population", "This report includes all members who completed at least one successful outbound coaching call during an active enrollment period within the reporting window. Engagement/outreach calls and calls outside of enrollment are excluded. Members without qualifying coaching activity are excluded from all metrics.")
dd_entry(dd1_tf, "Wellbeing Topics", "Each coaching call is assigned a single discussion topic via a tiered classification system: (1) Coach’s direct form selection, (2) keyword inference from coach notes, (3) call type mapping, or (4) most recent prior topic within 180 days. Calls with no determinable topic are classified as ‘General.’")
dd_entry(dd1_tf, "Goal Domains", "Coaching goals are categorized into 11 wellbeing domains: Gaps in Care, Exercise, Nutrition, Weight Management, Tobacco Cessation, Mental/Behavioral Health, Stress Management, Condition Management, Financial, Social, and Spiritual. Goals that do not match a recognized domain are reported as ‘Other.’")

dd_section(dd1_tf, "Goal Status Definitions")
dd_entry(dd1_tf, "Not Started", "A goal created in the system but the member has not yet begun working toward it. Typically system-recommended goals from evidence-based goal libraries awaiting member activation.")
dd_entry(dd1_tf, "In Progress", "A goal the member is actively working on with their coach.")
dd_entry(dd1_tf, "Completed", "A goal that was achieved and formally closed by the coach or system.")
dd_entry(dd1_tf, "Withdrawn", "A goal that was abandoned, refused by the member, or determined to be no longer applicable.")
dd_entry(dd1_tf, "Completion Rate", "Goals completed divided by total goals (Completed + In Progress + Not Started + Withdrawn), expressed as a percentage.")

dd_section(dd1_tf, "Tobacco Coaching Focus")
dd_entry(dd1_tf, "Tobacco Participants", "Members who had ‘Tobacco’ selected as their coaching discussion topic or whose call type was ‘Tobacco’ at any point in their coaching history, and who had at least one call in the reporting period.")
dd_entry(dd1_tf, "Active Tobacco Participants", "Subset of Tobacco Participants who additionally have a formal Tobacco Cessation goal with status In Progress or Not Started in the coaching platform.")
dd_entry(dd1_tf, "Goal Coverage Gap", "The number and percentage of Tobacco Participants who discussed tobacco during coaching but do not have a formal Tobacco Cessation goal being tracked. These members may have goals in other domains.")


# ===== SLIDE 7: Data Dictionary (Page 2) =====
slide = prs.slides.add_slide(BLANK)
add_title(slide, "Data Dictionary (continued)", Inches(0.4), Inches(0.15), size=16)

line = slide.shapes.add_connector(1, Inches(0.4), Inches(0.5), Inches(9.6), Inches(0.5))
line.line.color.rgb = COLOR_ACCENT
line.line.width = Pt(3)

dd2 = slide.shapes.add_textbox(Inches(0.4), Inches(0.6), Inches(9.2), Inches(7.0))
dd2_tf = dd2.text_frame
dd2_tf.word_wrap = True

p0 = dd2_tf.paragraphs[0]
p0.text = "Timeframes & Comparisons"
p0.font.name = FONT_BODY
p0.font.size = Pt(11)
p0.font.bold = True
p0.font.color.rgb = COLOR_PRIMARY

dd_entry(dd2_tf, "L12M (Last 12 Months)", "Trailing 12 calendar months ending with the report month. Provides a comprehensive longitudinal view that smooths seasonal variation and short-term fluctuations.")
dd_entry(dd2_tf, "Current Month", "The most recent complete calendar month in the reporting period. Shown alongside the prior month for trend comparison.")
dd_entry(dd2_tf, "MoM Change", "Month-over-Month change. The net numerical difference between the current reporting month and the immediately prior month.")
dd_entry(dd2_tf, "90-Day Rolling Average (90D Avg)", "Total members or goals in the most recent 90-day window divided by 3, representing the average monthly run rate. Smooths week-to-week variability.")
dd_entry(dd2_tf, "Prior 90-Day Average (Prior 90D Avg)", "Same calculation applied to the 90-day window immediately preceding the current 90 days. Used as the baseline for 90D trend comparison.")
dd_entry(dd2_tf, "90D Delta (90D Δ)", "The difference between the current 90-Day Average and the Prior 90-Day Average. Indicates whether the metric is trending up or down over a 6-month horizon.")

dd_section(dd2_tf, "Data Sources & Methodology")
dd_entry(dd2_tf, "Coaching Calls", "ENT_WH.MEMBER_CALL_DATA – successful outbound calls filtered to billable or interaction-eligible call types per ENT_WH.CALLTYPE_XREF_VW. Restricted to calls during active enrollment (LEVEL_NAME = Enrolled) and excludes Engagement call types (outreach/scheduling). One call per member per day (most recent kept).")
dd_entry(dd2_tf, "Discussion Topics", "ENT_WH.COACH_NOTES_WORKFLOW (questions 502533, 502758 for topic selection; 502534, 502833 for goal/detail text) and ENT_WH.COACH_NOTES_WORKFLOW_DM (questions 502599, 502616 for disease management).")
dd_entry(dd2_tf, "Goals", "SCP.AH_MEMBER_ACTION – ActionType 2 (Coach-Created) and ActionType 3 (System-Recommended). Status reflects the current snapshot as of the data refresh date, not historical progression.")
dd_entry(dd2_tf, "Member Identity", "All member counts use CURRENTGUID as the canonical identifier to correctly handle GUID shifts. GUID resolution performed via ENT_WH.ELIGMEMBER.")
dd_entry(dd2_tf, "Enrollment Bridge", "ENT_WH.COACHING_ENROLLMENT_MODEL links coaching account IDs to member GUIDs and CURRENTGUIDs for the goals section.")

dd_section(dd2_tf, "Refresh & Delivery")
dd_entry(dd2_tf, "Data Refresh", "All underlying data tables are refreshed on the 1st of each month. Metrics reflect the state of the data as of the most recent refresh.")
dd_entry(dd2_tf, "Report Generation", "This report is generated automatically from Vertica data tables (Carefirst_Sandbox schema) using a parameterized Python pipeline. It can be regenerated for any customer or date range on demand.")
dd_entry(dd2_tf, "Contact", "For questions about methodology, data definitions, or custom report requests, contact the Client Reporting Services team.")


# --- Save ---
period_str = f"{start_date.replace('-', '')}_{end_date.replace('-', '')}"
output_filename = f"coaching_report_{customer_id}_{period_str}.pptx"
output_path = os.path.join(SCRIPT_DIR, output_filename)
prs.save(output_path)
log.info(f"Report saved: {output_path}")
print(f"\nDone! {output_path}")
